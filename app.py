import streamlit as st
import os
import tempfile
import anthropic
from streamlit_mermaid import st_mermaid
from PyPDF2 import PdfReader
from pathlib import Path
import time
from io import BytesIO
from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle

# This MUST be the very first Streamlit command
st.set_page_config(
    page_title="PitchMe",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Constants for prompts
STORY_PROMPT = """
# Role
You are a Startup Pitch Evaluator with expertise in storytelling and pitch deck evaluation.

# Task
Analyze the following pitch deck and identify the story it tells. Focus on the customer's problem, the solution, and their journey.

If there is no clear story, explain this and suggest a customer story for the pitch.

# Output Format
Provide your response in two sections with clear markdown headers (## for main sections, ### for subsections):

## 📝 Story Analysis
A brief summary of the story told in the pitch (or lack thereof). Use bullet points for clarity.

## 🔍 Storytelling Recommendations
Creative recommendations for improving the storytelling based on three approaches. Format each approach as a subsection (###) with clear recommendations:

### The Hero's Journey
- **Strengths**: "This is good because..." 
- **Improvements**: "I suggest you can improve this by..."

### The Customer's Tale
- **Strengths**: "This is good because..." 
- **Improvements**: "I suggest you can improve this by..."

### The Industry's Point of View
- **Strengths**: "This is good because..." 
- **Improvements**: "I suggest you can improve this by..."

Use bold text for important points, create tables for comparing approaches, and utilize emojis to make the content more visually engaging.

# Pitch Deck Content
{pitch_deck_text}
"""

STARTUP_STAGE_PROMPT = """
# Role
You are a Startup Stage Evaluator with expertise in identifying a venture's growth stage.

# Task
Analyze the following pitch deck and identify which of these startup stages the venture is at:
- Ideation: Conceptualization of core idea, defining business model, target audience, and potential solutions
- Minimum Viable Category (MVC): Identified a unique market category to potentially dominate, outlined problem, audience, and market type
- Initial Product Release (IPR): Launch of first product/service, often as beta for early adopters
- Minimum Viable Product (MVP): Product with just enough features to satisfy early customers and provide feedback
- Minimum Viable Repeatability (MVR): Consistently delivering to target market, demonstrating repeatable business model and scalability potential

# Output Format
Provide your response using clear markdown formatting:

## 🚀 Current Stage
Create a visual indicator showing where the startup is on the journey, like:
Ideation → [MVC] → IPR → MVP → MVR

Explain which stage the startup is at and justify with specific evidence from the pitch deck. Use bullet points to list the evidence.

## 📋 Stage Definition
Provide the definition of the identified stage with clear formatting.

## 🔮 Next Stage Planning
Create a visual roadmap showing how to get to the next stage. Use a table format with these columns:
| Current Status | Next Milestone | Action Items |
Include 3-5 key actions the startup should take to reach the next stage.

Use bold text for important points and emojis to make the content visually engaging.

# Pitch Deck Content
{pitch_deck_text}
"""

MARKET_ENTRY_PROMPT = """
# Role
You are a Market Strategy Expert who specializes in evaluating startup market entry approaches.

# Task
Analyze the following pitch deck and determine:
1. If the startup has clearly identified their most critical customer segment
2. Whether their strategy aligns more with Blue Ocean Strategy or Red Ocean Strategy

Blue Ocean Strategy:
- Creates new, uncontested market space
- Makes competition irrelevant
- Creates and captures new demand
- Breaks the value-cost trade-off
- Achieves differentiation and low cost simultaneously

Red Ocean Strategy:
- Competes in existing market space
- Beats the competition
- Exploits existing demand
- Makes the value-cost trade-off
- Chooses between differentiation or low cost

# Output Format
Provide your response using clear markdown formatting with visual elements:

## 🎯 Customer Segment Analysis
Create a target diagram using mermaid syntax to represent how well they've identified their critical customer segment.

Evaluate how well they've identified their critical customer segment, using bullet points for clarity and bold text for key insights.

## 🌊 Strategy Classification
Generate a mermaid diagram that visually represents the position on the Blue Ocean vs Red Ocean spectrum. The diagram should be a horizontal flow with nodes indicating "Red Ocean" and "Blue Ocean" and an indicator node positioned accordingly.

Determine if they're using Blue Ocean or Red Ocean strategy with evidence. Use a comparison table:

| Blue Ocean Indicators | Red Ocean Indicators |
| --------------------- | -------------------- |
| (List evidence) | (List evidence) |

## 📈 Market Entry Recommendations
Provide recommendations to sharpen their customer segment focus and strengthen their chosen strategy. Format as:

### Strengths:
- "This is good because..." (bullet points)

### Improvements:
- "I suggest you can improve this by..." (bullet points)

Use emojis, bold text, and clear formatting to make the content visually engaging.

# Pitch Deck Content
{pitch_deck_text}
"""

BUSINESS_MODEL_PROMPT = """
# Role
You are a Business Model Expert specializing in startup evaluation.

# Task
Evaluate the following pitch deck using the Business Model Canvas framework. Score each element from 0-10 (0 = not addressed, 10 = excellent).

# Output Format
Create a visually appealing report with clear markdown formatting:

## 💼 Business Model Canvas Evaluation

For each Business Model Canvas element, create a subsection with the following format:

### [Element Name] 📊 Score: [X/10]

#### Evidence:
> Quote relevant text from the pitch deck (in blockquote format)

#### Strengths:
- ✅ "This is good because..." (bullet points)

#### Areas for Improvement:
- 🔄 "I suggest you can improve this by..." (bullet points)
OR
- ❓ "Not addressed. Go out and talk to experts! Try contacting... and ask them..."
OR
- 🔍 "Not addressed. How does your competitor's business or financial model address..."

Include these elements in your evaluation:
- Customer Segments
- Value Propositions
- Channels
- Revenue Streams
- Customer Relationships
- Key Activities
- Key Resources
- Key Partners
- Cost Structure

End with a visual summary table showing scores for all elements:

| Element | Score | Key Insight |
| ------- | ----- | ----------- |
| [Element] | [Score] | [Brief comment] |

Also create a radar chart representation using mermaid syntax to visualize the scores across all elements.

# Pitch Deck Content
{pitch_deck_text}
"""

EXPERT_PANEL_PROMPT = """
# Role
You are a Panel Moderator hosting a group of startup experts.

# Task
Simulate feedback from a panel of 5 experts reviewing the following pitch deck:
- Product Expert: Evaluates functionality, usability, design, and product quality
- Revenue Expert: Analyzes monetization strategies, sales, marketing, and revenue potential
- Team Expert: Assesses skills, experience, and cohesiveness of the startup team
- System Expert: Examines operations, processes, and systems for scalability
- Subject Matter Expert: Provides industry-specific insights and market knowledge

# Output Format
Create a visually engaging panel discussion report with clear markdown formatting:

## 👥 Expert Panel Feedback

For each expert, create a profile and feedback section:

### 👨‍💼 [Expert Title]
**Focus Areas**: [key evaluation criteria]

#### Key Observations:
> Quote relevant text from the pitch deck (in blockquote format)

#### Feedback:
- ✅ **Strengths**: "This is good because..." (bullet points)
- 🔄 **Suggestions**: "I suggest you can improve this by..." (bullet points)
- ❓ **Questions**: "The expert would ask..." (if applicable)

End with a panel summary showing the overall consensus, areas of agreement, and any conflicting viewpoints. Create a table showing each expert's key recommendation:

| Expert | Key Recommendation | Priority Level |
| ------ | ------------------ | -------------- |
| [Expert] | [Recommendation] | [High/Medium/Low] |

Use emojis, bold text, and clear formatting to make the content visually engaging.

# Pitch Deck Content
{pitch_deck_text}
"""

OVERALL_FEEDBACK_PROMPT = """
# Role
You are a Startup Mentor with expertise in pitch deck evaluation and fostering a learning mindset.

# Task
Provide comprehensive feedback on the following pitch deck.

# Output Format
Create a visually engaging executive summary with clear markdown formatting:

## 📝 Executive Summary

### ✨ Strengths
Create a visual scorecard for 3-5 key strengths of the pitch deck. For each strength:
- **[Strength Title]**: Detailed explanation with specific examples from the pitch deck
- **Impact**: Why this matters for investors and customers
- **Leverage Point**: How to maximize this strength

### 🔍 Areas for Improvement
Create a prioritized list of 3-5 specific areas where the pitch could be enhanced:
- **[Area Title]** (Priority: High/Medium/Low)
  - **Current State**: What the pitch currently shows
  - **Desired State**: What would make it more compelling
  - **Gap Analysis**: What's missing and why it matters

### 🚀 Action Plan
Create a table with 3-5 concrete actions:

| Action Item | Expected Impact | Difficulty | Timeline |
| ----------- | --------------- | ---------- | -------- |
| [Action] | [Impact] | [Easy/Medium/Hard] | [Timeframe] |

### 💭 Motivational Closing
A paragraph that encourages the team to view feedback as an opportunity for growth, emphasizing the learning mindset. Use metaphors and inspirational language that connects to the startup's mission.

Use emojis, bold text, tables, and clear formatting to make the content visually engaging.

# Pitch Deck Content
{pitch_deck_text}
"""

DESIGN_ANALYSIS_PROMPT = """
# Role
You are a Design and Visual Communication Expert specializing in evaluating pitch deck visuals and design elements.

# Task
Analyze the presentation text provided below and infer what design and visual elements are likely present based on the content. Even though you can't directly see the images, you can make educated evaluations based on:

1. References to visual elements (charts, graphs, images, diagrams)
2. The structure and flow of information
3. Mentions of branding, colors, or visual elements
4. Layout descriptions or implied formatting

# Output Format
Create a visually engaging design analysis with clear markdown formatting:

## 🎨 Design Elements Analysis
Create a visual checklist of design elements likely present in the pitch deck:

- [ ] Professional color scheme
- [ ] Consistent typography
- [ ] High-quality images
- [ ] Effective charts/graphs
- [ ] Clear slide layouts
- [ ] Visual hierarchy
- [ ] Branded elements

For each element detected, change [ ] to [x] and provide evidence from the content.

## 🔍 Visual Branding Evaluation
Create a mock brand guideline based on the inferred elements:
- **Colors**: Likely palette (based on any color mentions)
- **Typography**: Inferred font choices and hierarchy
- **Imagery**: Types of visuals mentioned
- **Layout**: Structure and organization patterns

## 💡 Design Recommendations
Organize recommendations by category:

### Slide Layouts
- ✅ **Strengths**: "This appears well-designed because..." (based on the content)
- 🔄 **Improvements**: "Consider improving this by..."

### Color Scheme
- ✅ **Strengths**: "This appears well-designed because..."
- 🔄 **Improvements**: "Consider improving this by..."

(Repeat for Typography, Charts/Diagrams, Image Selection, Visual Storytelling)

End with a visual "before/after" concept using mermaid syntax to illustrate key improvements.

Use emojis, bold text, and clear formatting to make the content visually engaging.

# Pitch Deck Content
{pitch_deck_text}
"""

# Add CSS for styling, with dark mode support
def add_custom_css():
    st.markdown("""
    <style>
        /* Basic styling */
        h1, h2, h3 { margin-bottom: 1rem; }
        p, li, div { line-height: 1.6; }
        
        /* Sidebar styling */
        [data-testid="stSidebar"] {
            background-color: #1a365d;
            padding-top: 1rem;
        }
        [data-testid="stSidebar"] p, 
        [data-testid="stSidebar"] h1, 
        [data-testid="stSidebar"] h2, 
        [data-testid="stSidebar"] h3, 
        [data-testid="stSidebar"] h4, 
        [data-testid="stSidebar"] span,
        [data-testid="stSidebar"] a,
        [data-testid="stSidebar"] .stMarkdown,
        [data-testid="stSidebar"] div,
        [data-testid="stSidebar"] label,
        [data-testid="stSidebar"] li,
        [data-testid="stSidebar"] .streamlit-expanderHeader,
        [data-testid="stSidebar"] .streamlit-expanderContent {
            color: #ffffff !important;
        }
        [data-testid="stSidebar"] h1 {
            font-size: 1.8rem;
            margin-bottom: 0.8rem;
            padding-bottom: 0.8rem;
            border-bottom: 1px solid rgba(255,255,255,0.2);
        }
        
        /* Tabs styling */
        .stTabs [data-baseweb="tab-list"] {
            gap: 0px;
            border-bottom: 1px solid #e2e8f0;
        }
        .stTabs [data-baseweb="tab"] {
            border-radius: 4px 4px 0px 0px;
            padding: 10px 16px;
            font-weight: 500;
            border: 1px solid #e2e8f0;
            border-bottom: none;
            margin-right: 4px;
        }
        
        /* Button styling */
        .stButton button {
            background-color: #2a70ba !important;
            color: white !important;
            border: none !important;
            font-weight: bold !important;
        }
        .stButton button:hover {
            background-color: #1a5ba6 !important;
        }
        
        /* Make blockquotes stand out */
        blockquote {
            border-left: 5px solid #2a70ba;
            padding: 10px 15px;
            margin: 10px 0;
        }
        
        /* Table styling */
        table {
            border-collapse: collapse;
            width: 100%;
            margin: 16px 0;
        }
        
        /* Responsive design adjustments */
        @media (max-width: 768px) {
            .row-widget.stButton {
                width: 100%;
            }
            
            /* Make tables scrollable on mobile */
            .stTable {
                overflow-x: auto;
                display: block;
            }
            
            /* Better spacing for mobile */
            .block-container {
                padding-left: 1rem;
                padding-right: 1rem;
            }
        }
        
        /* Dark mode support */
        @media (prefers-color-scheme: dark) {
            .stApp {
                background-color: #0e1117;
            }
            
            p, li, span, div {
                color: #f9fafb !important;
            }
            
            h1, h2, h3, h4, h5, h6 {
                color: #e5e7eb !important;
            }
            
            .stTabs [data-baseweb="tab"] {
                background-color: #1e2530;
                color: #e5e7eb;
                border-color: #4b5563;
            }
            
            table {
                border-color: #4b5563;
            }
            
            th, td {
                border-color: #4b5563;
                color: #e5e7eb;
            }
            
            blockquote {
                background-color: #1e2530;
                color: #e5e7eb;
            }
        }
        
        /* Upload section styling */
        .upload-section {
            border-radius: 8px;
            padding: 20px;
            margin-bottom: 20px;
        }
        
        /* Features section styling */
        .features-section {
            border-radius: 8px;
            padding: 20px;
            margin-bottom: 20px;
        }
        
        /* Supported formats section styling */
        .formats-section {
            border-radius: 8px;
            padding: 20px;
            margin-bottom: 20px;
        }
        
        /* Dark mode styling for sections */
        @media (prefers-color-scheme: dark) {
            .upload-section, .features-section, .formats-section {
                background-color: #1e2530;
            }
        }
        
        /* Light mode styling for sections */
        @media (prefers-color-scheme: light) {
            .upload-section, .features-section, .formats-section {
                background-color: white;
                box-shadow: 0 1px 3px rgba(0, 0, 0, 0.1);
            }
        }
        
        /* Emoji sizing */
        em {
            font-style: normal;
        }
    </style>
    """, unsafe_allow_html=True)

add_custom_css()

# Function to display company logo
def get_company_logo():
    try:
        # Use the file path relative to the app.py file
        return st.image("ProtoBots_Logo_Circle.png", width=80)
    except Exception as e:
        st.error(f"Error loading logo: {str(e)}")
        # Fallback to text if image fails to load
        return st.markdown("<h3>ProtoBots.ai</h3>", unsafe_allow_html=True)

# Create Anthropic client
def get_anthropic_client():
    api_key = os.environ.get("ANTHROPIC_API_KEY")
    if not api_key:
        try:
            api_key = st.secrets["ANTHROPIC_API_KEY"]
        except Exception as e:
            st.error(f"Error accessing secrets: {str(e)}")
            st.stop()
    
    # Try different initialization methods for compatibility
    try:
        return anthropic.Anthropic(api_key=api_key)
    except Exception as e:
        try:
            return anthropic.Client(api_key=api_key)
        except Exception as e:
            st.error(f"Could not initialize Anthropic client: {str(e)}")
            st.stop()

client = get_anthropic_client()

# Function to call Claude API
def call_claude_api(prompt, max_tokens=4000):
    try:
        # Try newer API first
        if hasattr(client, 'messages'):
            message = client.messages.create(
                model="claude-3-5-sonnet-20240620",
                max_tokens=max_tokens,
                messages=[{"role": "user", "content": prompt}]
            )
            return message.content[0].text
        # Fall back to older API
        else:
            response = client.completion(
                prompt=f"\n\nHuman: {prompt}\n\nAssistant:",
                model="claude-3-5-sonnet-20240620",
                max_tokens_to_sample=max_tokens,
                stop_sequences=["\n\nHuman:"]
            )
            return response.completion
    except Exception as e:
        st.error(f"Error calling Claude API: {str(e)}")
        return None

# Extract text from various file formats
def extract_text_from_file(uploaded_file):
    file_extension = uploaded_file.name.split('.')[-1].lower()
    if file_extension == 'pdf':
        return extract_text_from_pdf(uploaded_file)
    elif file_extension in ['ppt', 'pptx']:
        return extract_text_from_pptx(uploaded_file)
    elif file_extension in ['doc', 'docx']:
        return extract_text_from_docx(uploaded_file)
    else:
        st.error(f"Unsupported file format: .{file_extension}")
        return None

# Extract text from PDF
def extract_text_from_pdf(pdf_file):
    temp_dir = tempfile.TemporaryDirectory()
    temp_path = Path(temp_dir.name) / "pitch_deck.pdf"
    with open(temp_path, "wb") as f:
        f.write(pdf_file.getvalue())
    pdf_reader = PdfReader(temp_path)
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text() + "\n\n"
    temp_dir.cleanup()
    return text

# Extract text from PowerPoint
def extract_text_from_pptx(pptx_file):
    try:
        import pptx
        temp_dir = tempfile.TemporaryDirectory()
        temp_path = Path(temp_dir.name) / "pitch_deck.pptx"
        with open(temp_path, "wb") as f:
            f.write(pptx_file.getvalue())
        presentation = pptx.Presentation(temp_path)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            text += "\n\n"
        temp_dir.cleanup()
        return text
    except ImportError:
        st.error("PowerPoint processing library not available. Please install python-pptx.")
        return None

# Extract text from Word document
def extract_text_from_docx(docx_file):
    try:
        import docx
        temp_dir = tempfile.TemporaryDirectory()
        temp_path = Path(temp_dir.name) / "pitch_deck.docx"
        with open(temp_path, "wb") as f:
            f.write(docx_file.getvalue())
        doc = docx.Document(temp_path)
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
        temp_dir.cleanup()
        return text
    except ImportError:
        st.error("Word processing library not available. Please install python-docx.")
        return None

# Function to export evaluation results as a PDF
def export_results_to_pdf(results):
    try:
        buffer = BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=letter)
        styles = getSampleStyleSheet()
        
        # Add custom style for headers
        styles.add(ParagraphStyle(
            name='CustomHeading1',
            parent=styles['Heading1'],
            fontSize=16,
            spaceAfter=12
        ))
        
        content = []
        
        # Add title
        content.append(Paragraph("PitchMe Analysis Report", styles['CustomHeading1']))
        content.append(Spacer(1, 12))
        
        # Add each section
        for section, section_content in results.items():
            # Format section title
            formatted_title = " ".join(word.capitalize() for word in section.replace("_", " ").split())
            content.append(Paragraph(formatted_title, styles['Heading2']))
            content.append(Spacer(1, 6))
            
            # Add plain text content - simple but reliable
            paragraphs = section_content.split('\n\n')
            for para in paragraphs:
                if para.strip():
                    content.append(Paragraph(para.replace("<", "&lt;").replace(">", "&gt;"), styles['Normal']))
                    content.append(Spacer(1, 6))
            
            content.append(Spacer(1, 12))
        
        # Build the PDF
        doc.build(content)
        buffer.seek(0)
        return buffer.getvalue()
    except Exception as e:
        st.error(f"Error creating PDF: {str(e)}")
        # Return a simple error message as PDF
        return b"Could not generate PDF report. See error in application."

# Function to evaluate the pitch deck
def evaluate_pitch_deck(pitch_deck_text, analyze_design=False):
    results = {}
    progress_bar = st.progress(0)
    status_text = st.empty()
    # Determine how many analyses we'll run
    total_analyses = 6  # Base analyses including business model and expert panel
    if analyze_design:
        total_analyses += 1
    progress_step = 100 / total_analyses
    current_progress = 0

    # Story Analysis
    status_text.text("Analyzing story elements...")
    story_prompt = STORY_PROMPT.format(pitch_deck_text=pitch_deck_text)
    story_analysis = call_claude_api(story_prompt)
    if story_analysis:
        results["story"] = story_analysis
        current_progress += progress_step
        progress_bar.progress(int(current_progress))
    else:
        st.error("Failed to analyze story elements.")
        return None

    # Startup Stage
    status_text.text("Identifying startup stage...")
    stage_prompt = STARTUP_STAGE_PROMPT.format(pitch_deck_text=pitch_deck_text)
    stage_analysis = call_claude_api(stage_prompt)
    if stage_analysis:
        results["startup_stage"] = stage_analysis
        current_progress += progress_step
        progress_bar.progress(int(current_progress))
    else:
        st.error("Failed to identify startup stage.")
        return None

    # Market Entry
    status_text.text("Evaluating market entry strategy...")
    market_prompt = MARKET_ENTRY_PROMPT.format(pitch_deck_text=pitch_deck_text)
    market_analysis = call_claude_api(market_prompt)
    if market_analysis:
        results["market_entry"] = market_analysis
        current_progress += progress_step
        progress_bar.progress(int(current_progress))
    else:
        st.error("Failed to evaluate market entry strategy.")
        return None

    # Business Model
    status_text.text("Analyzing business model...")
    business_prompt = BUSINESS_MODEL_PROMPT.format(pitch_deck_text=pitch_deck_text)
    business_analysis = call_claude_api(business_prompt, max_tokens=6000)
    if business_analysis:
        results["business_model"] = business_analysis
        current_progress += progress_step
        progress_bar.progress(int(current_progress))
    else:
        st.error("Failed to analyze business model.")
        return None

    # Expert Panel
    status_text.text("Gathering expert panel feedback...")
    expert_prompt = EXPERT_PANEL_PROMPT.format(pitch_deck_text=pitch_deck_text)
    expert_analysis = call_claude_api(expert_prompt, max_tokens=6000)
    if expert_analysis:
        results["expert_panel"] = expert_analysis
        current_progress += progress_step
        progress_bar.progress(int(current_progress))
    else:
        st.error("Failed to gather expert panel feedback.")
        return None

    # Design Analysis (optional)
    if analyze_design:
        status_text.text("Analyzing design elements...")
        design_prompt = DESIGN_ANALYSIS_PROMPT.format(pitch_deck_text=pitch_deck_text)
        design_analysis = call_claude_api(design_prompt)
        if design_analysis:
            results["design"] = design_analysis
            current_progress += progress_step
            progress_bar.progress(int(current_progress))

    # Overall Feedback
    status_text.text("Generating overall feedback...")
    feedback_prompt = OVERALL_FEEDBACK_PROMPT.format(pitch_deck_text=pitch_deck_text)
    overall_feedback = call_claude_api(feedback_prompt)
    if overall_feedback:
        results["overall_feedback"] = overall_feedback
        current_progress += progress_step
        progress_bar.progress(100)  # Ensure we reach 100%
        status_text.text("Analysis complete!")
    else:
        st.error("Failed to generate overall feedback.")
        return None

    return results

# Function to display evaluation results in tabs
def display_evaluation_results(results):
    st.title("Pitch Deck Evaluation")
    # Define all possible tabs and their keys
    tab_definitions = [
        {"label": "📖 Story Analysis", "key": "story"},
        {"label": "🚀 Startup Stage", "key": "startup_stage"},
        {"label": "🎯 Market Entry", "key": "market_entry"},
        {"label": "💼 Business Model", "key": "business_model"},
        {"label": "👥 Expert Panel", "key": "expert_panel"},
        {"label": "🎨 Design Analysis", "key": "design"},
        {"label": "📝 Overall Feedback", "key": "overall_feedback"}
    ]
    # Filter to include only tabs with results
    available_tabs = [tab for tab in tab_definitions if tab["key"] in results]
    tab_labels = [tab["label"] for tab in available_tabs]
    # Create tabs
    tabs = st.tabs(tab_labels)
    # Populate each tab with content
    for i, tab in enumerate(available_tabs):
        with tabs[i]:
            content = results[tab["key"]]
            
            # Regular markdown content (without mermaid)
            markdown_parts = []
            
            # Extract mermaid diagrams
            mermaid_blocks = []
            in_mermaid = False
            current_mermaid = []
            
            for line in content.split('\n'):
                if line.strip() == "```mermaid":
                    in_mermaid = True
                elif line.strip() == "```" and in_mermaid:
                    in_mermaid = False
                    mermaid_blocks.append('\n'.join(current_mermaid))
                    current_mermaid = []
                    # Add a placeholder for the mermaid diagram
                    markdown_parts.append("[MERMAID_DIAGRAM_" + str(len(mermaid_blocks) - 1) + "]")
                elif in_mermaid:
                    current_mermaid.append(line)
                else:
                    markdown_parts.append(line)
            
            # Join the markdown parts into a single string
            markdown_content = '\n'.join(markdown_parts)
            
            # Split by mermaid diagram placeholders
            markdown_segments = markdown_content.split("[MERMAID_DIAGRAM_")
            
            # Render first segment
            if markdown_segments[0]:
                st.markdown(markdown_segments[0])
            
            # Render the rest with mermaid diagrams
            for i, segment in enumerate(markdown_segments[1:], 0):
                idx, rest = segment.split("]", 1)
                
                # Render the mermaid diagram
                st_mermaid(mermaid_blocks[int(idx)])
                
                # Render the rest of the markdown
                if rest:
                    st.markdown(rest)

def main():
    # Sidebar
    with st.sidebar:
        # Display company logo at the top
        get_company_logo()
        st.title("PitchMe")
        st.markdown("Upload your pitch deck to get expert evaluation.")
        with st.expander("ℹ️ About"):
            st.markdown("""
            This app uses AI to evaluate your pitch deck from multiple angles:
            - Storytelling effectiveness
            - Startup stage identification
            - Market entry strategy analysis
            - Business model canvas evaluation
            - Expert panel feedback
            - Design and visual elements (optional)
            - Overall recommendations
            
            Contact us at Support@ProtoBots.ai
            """)
        with st.expander("📋 How to use"):
            st.markdown("""
            1. Enter your startup's name (optional)
            2. Upload your pitch deck (PDF, PPT, PPTX, DOC, or DOCX)
            3. Select whether to analyze design elements
            4. Click "Evaluate Pitch Deck"
            5. Review the detailed analysis across various tabs
            6. Use the feedback to improve your pitch deck
            """)
        st.divider()
        st.markdown("<div style='text-align: center; font-size: 0.9rem; opacity: 0.8; margin-top: 20px;'>Made by ProtoBots.ai</div>", unsafe_allow_html=True)
    
    # Use container to dynamically update content without page refresh
    main_container = st.container()
    with main_container:
        if "evaluation_results" not in st.session_state:
            # Initial state - show upload form
            # Replace the blank banner with grey horizontal lines and title/subtitle
            st.markdown("<hr style='border: none; height: 2px; background: #ccc; box-shadow: 0 2px 2px -2px grey;'>", unsafe_allow_html=True)
            st.title("PitchMe")
            st.markdown("Get expert AI-powered feedback on your pitch deck to impress investors and secure funding.")
            st.markdown("<hr style='border: none; height: 2px; background: #ccc; box-shadow: 0 2px 2px -2px grey;'>", unsafe_allow_html=True)
            
            col1, col2 = st.columns([2, 1])
            with col1:
                st.markdown("<div class='upload-section'>", unsafe_allow_html=True)
                st.header("Upload Your Pitch Deck")
                startup_name = st.text_input("Startup Name (Optional)", "")
                uploaded_file = st.file_uploader(
                    "Upload your pitch deck", 
                    type=["pdf", "ppt", "pptx", "doc", "docx"]
                )
                analyze_design = st.checkbox("Also analyze design and visual elements", value=True)
                if uploaded_file is not None:
                    file_type = uploaded_file.name.split('.')[-1].lower()
                    st.write(f"File type detected: .{file_type}")
                    st.markdown("""<style>
                    div.stButton > button {
                        font-weight: bold !important;
                        color: white !important;
                    }
                    </style>""", unsafe_allow_html=True)
                    if st.button("Evaluate Pitch Deck", type="primary", use_container_width=True):
                        pitch_deck_text = extract_text_from_file(uploaded_file)
                        if not pitch_deck_text or len(pitch_deck_text) < 100:
                            st.error("Could not extract sufficient text from the file. Please make sure your file has textual content and not just images.")
                        else:
                            st.session_state.startup_name = startup_name
                            analysis_status = st.empty()
                            with analysis_status.container():
                                with st.spinner("Analyzing your pitch deck..."):
                                    results = evaluate_pitch_deck(pitch_deck_text, analyze_design)
                                    if results:
                                        st.session_state.evaluation_results = results
                                        st.success("Analysis complete! Displaying results...")
                                        time.sleep(1)
                                        main_container.empty()
                                        display_evaluation_results(results)
                                        # Add download button to export analysis as PDF
                                        pdf_bytes = export_results_to_pdf(st.session_state.evaluation_results)
                                        st.download_button(
                                            label="Export Analysis as PDF",
                                            data=pdf_bytes,
                                            file_name="PitchMe_Analysis.pdf",
                                            mime="application/pdf"
                                        )
                                        if st.button("Evaluate Another Pitch Deck", type="primary"):
                                            for key in list(st.session_state.keys()):
                                                del st.session_state[key]
                                            st.experimental_rerun()
                st.markdown("</div>", unsafe_allow_html=True)
            with col2:
                st.markdown("<div class='features-section'>", unsafe_allow_html=True)
                st.header("What You'll Get")
                st.markdown("""
                - 📖 **Story Analysis**
                - 🚀 **Startup Stage Identification**
                - 🎯 **Market Entry Strategy Assessment**
                - 💼 **Business Model Evaluation**
                - 👥 **Expert Panel Feedback**
                - 🎨 **Design Analysis** (optional)
                - 📝 **Actionable Recommendations**
                """)
                st.markdown("</div>", unsafe_allow_html=True)
                st.markdown("<div class='formats-section'>", unsafe_allow_html=True)
                st.header("Supported Formats")
                st.markdown("""
                We support multiple presentation formats:
                
                - **PDF** (.pdf)
                - **PowerPoint** (.ppt, .pptx)
                - **Word Documents** (.doc, .docx)
                """)
                st.markdown("</div>", unsafe_allow_html=True)
        else:
            display_evaluation_results(st.session_state.evaluation_results)
            if st.button("Evaluate Another Pitch Deck", type="primary"):
                for key in list(st.session_state.keys()):
                    del st.session_state[key]
                st.experimental_rerun()

if __name__ == "__main__":
    main()
